"""Generate the HealthMap pitch deck as a .pptx file."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Theme colours ─────────────────────────────────────────────────────
BG = RGBColor(0x0F, 0x17, 0x24)          # dark navy
ACCENT = RGBColor(0x2D, 0xD4, 0xBF)      # teal
ACCENT2 = RGBColor(0x38, 0xBD, 0xF8)     # sky blue
WHITE = RGBColor(0xE5, 0xE5, 0xE5)
MUTED = RGBColor(0xA1, 0xA1, 0xAA)
SURFACE = RGBColor(0x1A, 0x24, 0x35)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)
GREEN = RGBColor(0x22, 0xC5, 0x5E)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

BLANK_LAYOUT = prs.slide_layouts[6]  # blank


# ── Helpers ───────────────────────────────────────────────────────────

def set_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, *,
             font_size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT,
             font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return tf


def add_para(tf, text, *, font_size=18, color=WHITE, bold=False,
             space_before=Pt(6), font_name="Calibri", align=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.space_before = space_before
    p.alignment = align
    return p


def add_bullet(tf, text, *, font_size=16, color=WHITE, bold=False,
               indent=Inches(0.3)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.space_before = Pt(4)
    p.level = 0
    pPr = p._pPr
    if pPr is None:
        from pptx.oxml.ns import qn
        pPr = p._p.get_or_add_pPr()
    pPr.set("marL", str(int(indent)))
    return p


def add_rounded_rect(slide, left, top, width, height, fill_color=SURFACE):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_accent_line(slide, left, top, width):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    return shape


def slide_number(slide, num):
    add_text(slide, SLIDE_W - Inches(1), SLIDE_H - Inches(0.5),
             Inches(0.8), Inches(0.4), str(num),
             font_size=11, color=MUTED, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s1)

# Accent line at top
add_accent_line(s1, Inches(0), Inches(0), SLIDE_W)

# Team badge
badge = add_rounded_rect(s1, Inches(9.5), Inches(1.0), Inches(3.0), Inches(0.55), SURFACE)
add_text(s1, Inches(9.5), Inches(1.05), Inches(3.0), Inches(0.5),
         "TEAM  3devs1git", font_size=16, color=ACCENT, bold=True,
         align=PP_ALIGN.CENTER)

# Main title
add_text(s1, Inches(1.0), Inches(2.0), Inches(11), Inches(1.2),
         "HealthMap", font_size=60, color=WHITE, bold=True)
add_text(s1, Inches(1.0), Inches(3.2), Inches(11), Inches(0.8),
         "AI-Powered Healthcare Data Harmonization", font_size=28, color=ACCENT)
add_text(s1, Inches(1.0), Inches(4.2), Inches(10), Inches(0.6),
         "Automatically map heterogeneous medical files into a unified, queryable schema.",
         font_size=18, color=MUTED)

# Bottom bar
add_rounded_rect(s1, Inches(1.0), Inches(5.8), Inches(11.3), Inches(0.7), SURFACE)
add_text(s1, Inches(1.3), Inches(5.85), Inches(10.7), Inches(0.6),
         "START Hack 2026  ·  On-Premises  ·  Privacy-First  ·  Zero Cloud Dependencies",
         font_size=14, color=MUTED, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ══════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s2)
add_accent_line(s2, Inches(0), Inches(0), SLIDE_W)
slide_number(s2, 2)

add_text(s2, Inches(1.0), Inches(0.5), Inches(11), Inches(0.7),
         "The Problem", font_size=36, color=WHITE, bold=True)
add_accent_line(s2, Inches(1.0), Inches(1.2), Inches(2.0))

# Pain points as cards
problems = [
    ("Fragmented Data", "Labs, devices, nursing, medications —\neach system exports in its own format\n(CSV, XLSX, PDF, different column names).", ACCENT2),
    ("Manual Mapping", "Trained staff spend hours matching\ncolumns to the unified schema.\nError-prone and doesn't scale.", ORANGE),
    ("No Unified View", "Without harmonized data, clinicians\ncannot see a complete patient picture\nacross departments.", RED),
]

for i, (title, desc, color) in enumerate(problems):
    left = Inches(1.0 + i * 3.9)
    add_rounded_rect(s2, left, Inches(1.8), Inches(3.5), Inches(3.5), SURFACE)
    # Colour bar at top of card
    bar = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(1.8), Inches(3.5), Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    add_text(s2, left + Inches(0.3), Inches(2.1), Inches(2.9), Inches(0.5),
             title, font_size=20, color=color, bold=True)
    add_text(s2, left + Inches(0.3), Inches(2.7), Inches(2.9), Inches(2.2),
             desc, font_size=15, color=MUTED)

# Bottom quote
add_text(s2, Inches(1.0), Inches(5.8), Inches(11), Inches(0.8),
         '"Healthcare generates 30% of the world\'s data, yet 97% of hospital data goes unused." — RBC Capital Markets',
         font_size=13, color=MUTED, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 3 — Our Solution
# ══════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s3)
add_accent_line(s3, Inches(0), Inches(0), SLIDE_W)
slide_number(s3, 3)

add_text(s3, Inches(1.0), Inches(0.5), Inches(11), Inches(0.7),
         "Our Solution", font_size=36, color=WHITE, bold=True)
add_accent_line(s3, Inches(1.0), Inches(1.2), Inches(2.0))

# Pipeline stages
stages = [
    ("1", "Upload", "Drag & drop any\nCSV, XLSX, or PDF", ACCENT2),
    ("2", "Classify", "AI identifies the\ntarget table", ACCENT),
    ("3", "Map", "AI maps columns\nto DB schema", PURPLE),
    ("4", "Review", "Human verifies\n& corrects", ORANGE),
    ("5", "Import", "Bulk insert into\nPostgreSQL", GREEN),
]

for i, (num, title, desc, color) in enumerate(stages):
    left = Inches(0.6 + i * 2.5)
    # Circle with number
    circle = s3.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.65), Inches(1.8), Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(22)
    p.font.color.rgb = BG
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)

    # Arrow (except last)
    if i < len(stages) - 1:
        arrow_left = left + Inches(1.55)
        add_text(s3, arrow_left, Inches(1.85), Inches(0.8), Inches(0.6),
                 "→", font_size=24, color=MUTED, align=PP_ALIGN.CENTER)

    add_text(s3, left + Inches(0.15), Inches(2.7), Inches(1.7), Inches(0.4),
             title, font_size=18, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(s3, left + Inches(0.0), Inches(3.1), Inches(2.0), Inches(0.8),
             desc, font_size=14, color=MUTED, align=PP_ALIGN.CENTER)

# Key points below
add_rounded_rect(s3, Inches(1.0), Inches(4.3), Inches(11.3), Inches(2.5), SURFACE)
points_tf = add_text(s3, Inches(1.4), Inches(4.4), Inches(10.5), Inches(2.3),
                     "Privacy-first: only column names, types, and 5 sample values sent to LLM — never raw patient data",
                     font_size=15, color=WHITE)
add_bullet(points_tf, "Real-time progress: SSE streaming shows each pipeline stage live in the browser", font_size=15, color=WHITE)
add_bullet(points_tf, "Smart caching: SHA256-keyed results avoid repeat LLM calls for files with the same schema", font_size=15, color=WHITE)
add_bullet(points_tf, "Human-in-the-loop: every mapping can be reviewed and corrected before import", font_size=15, color=WHITE)
add_bullet(points_tf, "Row-level error recovery: bad rows are skipped, good rows are imported — no all-or-nothing failures", font_size=15, color=WHITE)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 4 — Architecture
# ══════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s4)
add_accent_line(s4, Inches(0), Inches(0), SLIDE_W)
slide_number(s4, 4)

add_text(s4, Inches(1.0), Inches(0.5), Inches(11), Inches(0.7),
         "Architecture", font_size=36, color=WHITE, bold=True)
add_accent_line(s4, Inches(1.0), Inches(1.2), Inches(2.0))

# Service boxes
services = [
    ("React 19 + TypeScript", "Frontend\n:4242", ACCENT2, Inches(0.8)),
    ("Go (Gin + GORM)", "REST API\n:8080", ACCENT, Inches(3.8)),
    ("Python (FastAPI)", "ML Pipeline\n:5001", PURPLE, Inches(6.8)),
    ("Ollama (qwen2.5)", "Local LLM\n:11434", ORANGE, Inches(9.8)),
]

for (tech, label, color, left) in services:
    box = add_rounded_rect(s4, left, Inches(1.8), Inches(2.6), Inches(1.8), SURFACE)
    # Top bar
    bar = s4.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(1.8), Inches(2.6), Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    add_text(s4, left + Inches(0.2), Inches(2.05), Inches(2.2), Inches(0.4),
             tech, font_size=13, color=color, bold=True)
    add_text(s4, left + Inches(0.2), Inches(2.5), Inches(2.2), Inches(0.8),
             label, font_size=16, color=MUTED)

# Arrows between services
for x in [Inches(3.4), Inches(6.4), Inches(9.4)]:
    add_text(s4, x, Inches(2.3), Inches(0.5), Inches(0.5),
             "→", font_size=28, color=MUTED, align=PP_ALIGN.CENTER)

# PostgreSQL at bottom center
pg_box = add_rounded_rect(s4, Inches(4.5), Inches(4.2), Inches(4.3), Inches(1.0), SURFACE)
bar = s4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(4.2), Inches(4.3), Pt(4))
bar.fill.solid()
bar.fill.fore_color.rgb = GREEN
bar.line.fill.background()
add_text(s4, Inches(4.7), Inches(4.35), Inches(3.9), Inches(0.35),
         "PostgreSQL 16", font_size=13, color=GREEN, bold=True)
add_text(s4, Inches(4.7), Inches(4.65), Inches(3.9), Inches(0.4),
         "8 target healthcare tables  ·  Auto-migrated  ·  Persistent storage",
         font_size=12, color=MUTED)

# Up arrows to PG
add_text(s4, Inches(5.5), Inches(3.7), Inches(0.5), Inches(0.5),
         "↓", font_size=24, color=MUTED, align=PP_ALIGN.CENTER)

# Key architecture points
add_rounded_rect(s4, Inches(0.8), Inches(5.6), Inches(11.7), Inches(1.5), SURFACE)
arch_tf = add_text(s4, Inches(1.2), Inches(5.7), Inches(11.0), Inches(1.3),
                   "100% on-premises — all 5 services run in Docker, zero external API calls",
                   font_size=14, color=WHITE)
add_bullet(arch_tf, "Local LLM (Ollama) — no data leaves the hospital network, GDPR/HIPAA compliant by design", font_size=14, color=WHITE)
add_bullet(arch_tf, "Docker Compose orchestration — single command deployment: docker compose up", font_size=14, color=WHITE)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 5 — Innovation & Tech
# ══════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s5)
add_accent_line(s5, Inches(0), Inches(0), SLIDE_W)
slide_number(s5, 5)

add_text(s5, Inches(1.0), Inches(0.5), Inches(11), Inches(0.7),
         "What Makes It Innovative", font_size=36, color=WHITE, bold=True)
add_accent_line(s5, Inches(1.0), Inches(1.2), Inches(2.0))

innovations = [
    ("Metadata-Only AI", "LLM sees column names, types, and\n5 sample values — never raw patient\nrecords. Privacy by architecture.", ACCENT, Inches(1.0)),
    ("Reference-Enhanced Prompts", "Real data samples from the target DB\nguide the LLM, boosting accuracy\nfor German & abbreviated headers.", PURPLE, Inches(4.55)),
    ("Batched Column Mapping", "Files with 100+ columns are split into\nbatches of 20, then deduped and\nvalidated against the actual schema.", ACCENT2, Inches(8.1)),
]

for (title, desc, color, left) in innovations:
    add_rounded_rect(s5, left, Inches(1.7), Inches(3.2), Inches(2.8), SURFACE)
    bar = s5.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(1.7), Inches(3.2), Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    add_text(s5, left + Inches(0.25), Inches(1.95), Inches(2.7), Inches(0.5),
             title, font_size=18, color=color, bold=True)
    add_text(s5, left + Inches(0.25), Inches(2.5), Inches(2.7), Inches(1.6),
             desc, font_size=14, color=MUTED)

# Tech stack bar
add_rounded_rect(s5, Inches(1.0), Inches(4.9), Inches(11.3), Inches(2.2), SURFACE)
add_text(s5, Inches(1.3), Inches(5.0), Inches(10.7), Inches(0.4),
         "Tech Stack", font_size=18, color=WHITE, bold=True)

techs = [
    ("Go + Gin", "REST API, concurrency"),
    ("Python + FastAPI", "ML pipeline, async"),
    ("React 19 + TypeScript", "Interactive UI"),
    ("Ollama (qwen2.5:1.5b)", "Local LLM inference"),
    ("PostgreSQL 16", "Relational storage"),
    ("Docker Compose", "One-command deploy"),
]

for i, (name, desc) in enumerate(techs):
    col = i % 3
    row = i // 3
    left = Inches(1.3 + col * 3.7)
    top = Inches(5.5 + row * 0.7)
    add_text(s5, left, top, Inches(1.6), Inches(0.35),
             name, font_size=14, color=ACCENT, bold=True)
    add_text(s5, left + Inches(1.7), top, Inches(1.8), Inches(0.35),
             desc, font_size=13, color=MUTED)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 6 — Real-World Impact & Cost
# ══════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s6)
add_accent_line(s6, Inches(0), Inches(0), SLIDE_W)
slide_number(s6, 6)

add_text(s6, Inches(1.0), Inches(0.5), Inches(11), Inches(0.7),
         "Real-World Impact", font_size=36, color=WHITE, bold=True)
add_accent_line(s6, Inches(1.0), Inches(1.2), Inches(2.0))

# Comparison table: Before vs After
add_rounded_rect(s6, Inches(1.0), Inches(1.7), Inches(5.2), Inches(3.0), SURFACE)
add_text(s6, Inches(1.3), Inches(1.85), Inches(4.6), Inches(0.4),
         "Before HealthMap", font_size=20, color=RED, bold=True)
before_tf = add_text(s6, Inches(1.3), Inches(2.3), Inches(4.6), Inches(2.2),
                     "Hours of manual column mapping per file", font_size=15, color=MUTED)
add_bullet(before_tf, "Trained staff required for each new format", font_size=15, color=MUTED)
add_bullet(before_tf, "Human errors go undetected until production", font_size=15, color=MUTED)
add_bullet(before_tf, "New departments = new mapping effort", font_size=15, color=MUTED)

add_rounded_rect(s6, Inches(6.8), Inches(1.7), Inches(5.5), Inches(3.0), SURFACE)
add_text(s6, Inches(7.1), Inches(1.85), Inches(5.0), Inches(0.4),
         "With HealthMap", font_size=20, color=GREEN, bold=True)
after_tf = add_text(s6, Inches(7.1), Inches(2.3), Inches(5.0), Inches(2.2),
                    "Minutes — upload, review AI suggestion, import", font_size=15, color=WHITE)
add_bullet(after_tf, "Any staff member can onboard new data", font_size=15, color=WHITE)
add_bullet(after_tf, "AI flags anomalies (high nulls, duplicates)", font_size=15, color=WHITE)
add_bullet(after_tf, "Schema-validated — invalid columns auto-rejected", font_size=15, color=WHITE)

# Arrow between
add_text(s6, Inches(6.0), Inches(2.6), Inches(1.0), Inches(0.6),
         "→", font_size=36, color=ACCENT, align=PP_ALIGN.CENTER)

# Cost section
add_rounded_rect(s6, Inches(1.0), Inches(5.0), Inches(11.3), Inches(2.2), SURFACE)
add_text(s6, Inches(1.3), Inches(5.1), Inches(10.7), Inches(0.4),
         "Cost & Deployment", font_size=18, color=WHITE, bold=True)

costs = [
    ("Infrastructure", "Any Linux server with 8GB RAM + GPU optional", ACCENT),
    ("LLM Cost", "Zero — local Ollama, no API fees, no per-token billing", GREEN),
    ("Licensing", "Open-source stack (Go, Python, React, PostgreSQL)", ACCENT2),
    ("Deployment", "Single docker compose up — under 5 minutes", PURPLE),
]

for i, (label, value, color) in enumerate(costs):
    col = i % 2
    row = i // 2
    left = Inches(1.3 + col * 5.5)
    top = Inches(5.55 + row * 0.65)
    add_text(s6, left, top, Inches(1.8), Inches(0.35),
             label, font_size=14, color=color, bold=True)
    add_text(s6, left + Inches(1.9), top, Inches(3.3), Inches(0.35),
             value, font_size=13, color=MUTED)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 7 — Why HealthMap
# ══════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s7)
add_accent_line(s7, Inches(0), Inches(0), SLIDE_W)
slide_number(s7, 7)

add_text(s7, Inches(1.0), Inches(0.5), Inches(11), Inches(0.7),
         "Why HealthMap?", font_size=36, color=WHITE, bold=True)
add_accent_line(s7, Inches(1.0), Inches(1.2), Inches(2.0))

reasons = [
    ("On-Premises", "Runs entirely inside\nthe hospital network.\nGDPR/HIPAA compliant\nby architecture.", ACCENT),
    ("Zero Data Leakage", "No cloud APIs, no data\nleaves the premises.\nPatient privacy is not\na feature — it's the default.", GREEN),
    ("Human-in-the-Loop", "AI suggests, humans\ndecide. Every mapping\ncan be reviewed and\ncorrected before import.", PURPLE),
    ("Format Agnostic", "CSV, XLSX, PDF, TSV —\nany tabular format,\nany column naming\nconvention.", ACCENT2),
    ("Extensible", "Add new target tables\nby updating the schema.\nNo code changes needed\nfor new data types.", ORANGE),
]

for i, (title, desc, color) in enumerate(reasons):
    left = Inches(0.5 + i * 2.5)
    add_rounded_rect(s7, left, Inches(1.7), Inches(2.2), Inches(3.2), SURFACE)
    bar = s7.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(1.7), Inches(2.2), Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    add_text(s7, left + Inches(0.2), Inches(1.95), Inches(1.8), Inches(0.5),
             title, font_size=16, color=color, bold=True)
    add_text(s7, left + Inches(0.2), Inches(2.5), Inches(1.8), Inches(2.0),
             desc, font_size=13, color=MUTED)

# Closing statement
add_rounded_rect(s7, Inches(2.5), Inches(5.3), Inches(8.3), Inches(1.5), SURFACE)
add_text(s7, Inches(2.8), Inches(5.45), Inches(7.7), Inches(0.5),
         "HealthMap is not a black box.", font_size=22, color=WHITE, bold=True,
         align=PP_ALIGN.CENTER)
add_text(s7, Inches(2.8), Inches(5.95), Inches(7.7), Inches(0.6),
         "It's a transparent, reviewable, on-premises AI assistant that makes\nhealthcare data usable — safely, quickly, and at near-zero cost.",
         font_size=16, color=MUTED, align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────────────
out_path = "/home/noahw/projects/StartHack2026/HealthMap_Pitch_3devs1git.pptx"
prs.save(out_path)
print(f"Saved to {out_path}")
